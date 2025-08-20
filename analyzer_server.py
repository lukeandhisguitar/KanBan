# ==============================================================================
# =================== 板块一 / SECTION 1: SETUP & HELPERS ======================
# ==============================================================================

# -*- coding: utf-8 -*-
import os
import io
import json
import logging
import urllib.request
from flask import Flask, request, jsonify
from flask_cors import CORS
from http import HTTPStatus
import dashscope

# --- 文件处理库 ---
import PyPDF2
import docx
try:
    from pptx import Presentation
except ImportError:
    print("警告：未安装 python-pptx 库，将无法解析 .pptx 文件。请运行 'pip install python-pptx'")
    Presentation = None

# ================== 配置区 ========================
# API_KEY_FILENAME 不再需要，因为我们将从请求头中获取API Key
ANALYSIS_MODEL = "qwen-max-longcontext"
MAX_FILE_UPLOADS = 5 # 最多同时上传的文件数量
# ==================================================

class BColors:
    OKGREEN = '\033[92m'
    FAIL = '\033[91m'
    WARNING = '\033[93m'
    ENDC = '\033[0m'

log = logging.getLogger('werkzeug')
log.setLevel(logging.ERROR)
app = Flask(__name__, static_url_path='', static_folder='.')
CORS(app)

# load_api_key_from_file 函数不再需要

def test_network_connection():
    """
    精简后的启动检测函数，只检测网络代理设置。
    """
    https_proxy = None # 在函数开始时初始化代理变量
    print("1. 正在自动检测系统代理...")
    try:
        proxies = urllib.request.getproxies()
        https_proxy = proxies.get('https', proxies.get('http'))
        if https_proxy:
            print(f"   - {BColors.OKGREEN}自动检测到系统代理: {https_proxy}{BColors.ENDC}")
            os.environ['HTTP_PROXY'] = https_proxy
            os.environ['HTTPS_PROXY'] = https_proxy
        else:
            print(f"   - {BColors.WARNING}未检测到系统代理，将进行直接连接。{BColors.ENDC}")
            os.environ.pop('HTTP_PROXY', None)
            os.environ.pop('HTTPS_PROXY', None)
        # 只要网络代理检测不抛出异常，就认为网络基本就绪
        return True, https_proxy
    except Exception as e:
        print(f"   - {BColors.FAIL}检测系统代理时发生错误: {e}{BColors.ENDC}")
        return False, None # 返回失败状态和无代理
  
def extract_text_from_file(file_stream, filename):
    text = ""
    file_extension = os.path.splitext(filename)[1].lower()
    try:
        if file_extension == '.txt':
            text = file_stream.read().decode('utf-8')
        elif file_extension == '.pdf':
            reader = PyPDF2.PdfReader(file_stream)
            for page in reader.pages:
                text += page.extract_text() or ""
        elif file_extension == '.docx':
            doc = docx.Document(file_stream)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_extension == '.pptx':
            if Presentation is None: return None, "解析PPTX失败：请先运行 'pip install python-pptx'"
            prs = Presentation(file_stream)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        else:
            return None, "不支持的文件类型"
        
        return text, None
    except Exception as e:
        return None, f"解析文件 '{filename}' 时出错: {str(e)}"

def call_qwen_api(model: str, prompt: str, api_key: str) -> dict:
    """
    修改后的API调用函数，现在需要传入API Key。
    """
    try:
        dashscope.api_key = api_key # 在每次调用时设置Key
        messages = [{'role': 'system', 'content': '你是一位顶级的商业分析师和投资银行家，极其擅长从商业计划书、信息备忘录等文档中提炼和构建结构化的深度分析报告。你的分析必须严格基于原文，逻辑严谨，格式清晰，绝对禁止杜撰任何信息。'},
                    {'role': 'user', 'content': prompt}]
        response = dashscope.Generation.call(model=model, messages=messages, result_format='message')
        if response.status_code == HTTPStatus.OK:
            finish_reason = response.output.choices[0].finish_reason
            content = response.output.choices[0].message.content
            if finish_reason == 'content_filter': return {"success": False, "error": "AI分析失败：内容触发安全策略。"}
            if content:
                cleaned_text = content.strip().replace('```json', '').replace('```', '').strip()
                return {"success": True, "content": cleaned_text}
            else: return {"success": False, "error": f"AI调用成功，但未能返回有效内容 (终止原因: {finish_reason})。"}
        else: return {"success": False, "error": f"请求失败: {response.code} - {response.message}"}
    except Exception as e: return {"success": False, "error": f"调用API时异常: {str(e)}"}

@app.route('/')
def index():
    """
    修改根路由，使其直接提供前端HTML文件。
    """
    return app.send_static_file('项目管理看板.html')

allowed_extensions = {'.pdf', '.docx', '.pptx', '.txt'}
def allowed_file(filename): return '.' in filename and os.path.splitext(filename)[1].lower() in allowed_extensions

# ==============================================================================
# =================== 板块二 / SECTION 2: CORE PROMPT & API ROUTE ===============
# ==============================================================================
@app.route('/analyze', methods=['POST'])
def analyze_files():
    # --- 关键修改：从请求头获取 API Key ---
    api_key = request.headers.get('X-Api-Key')
    if not api_key:
        return jsonify({"error": "请求头中缺少 API Key (X-Api-Key)"}), 401 # 401 Unauthorized

    if 'files' not in request.files:
        return jsonify({"error": "请求中未找到文件字段(files)"}), 400
    
    files = request.files.getlist('files')
    
    if not files or all(f.filename == '' for f in files):
        return jsonify({"error": "未选择任何文件"}), 400

    all_texts = []
    for file in files:
        if file and allowed_file(file.filename):
            file.stream.seek(0)
            file_stream = io.BytesIO(file.read())
            extracted_text, error = extract_text_from_file(file_stream, file.filename)
            
            if error:
                return jsonify({"error": error}), 400
            if extracted_text and extracted_text.strip():
                all_texts.append(extracted_text)

    if not all_texts:
        return jsonify({"error": "无法从任何文件中提取有效文本内容。请检查文件是否为空、已损坏或为不支持的图片格式。"}), 400

    full_text = "\n\n--- NEW DOCUMENT ---\n\n".join(all_texts)



      # --- 股权投资模板 (Equity Investment Template) - 基本已经是完美的模板了 ---
    equity_summary_template = """### 1. 核心痛点
  请深入分析公司所处行业在“旧时代”的现状：在公司的产品或服务出现之前，市场上的主流解决方案是什么？这些传统方法给客户带来了哪些具体的不便、高昂的成本或低下的效率？行业内是否存在长期难以逾越的技术或商业模式瓶颈？最后，清晰阐述该公司是如何通过其独特的技术、创新的商业模式或差异化的服务，精准地解决了这些核心痛点。

### 2. 产品与服务
  **必须使用Markdown表格** 来清晰展示公司的产品或服务矩阵。有几个核心产品/服务就写几行，确保每一项都描述清晰。
  ```
  | 产品/服务名称 | 核心功能与特性描述 | 主要应用场景与目标客户 |
  |---|---|---|
  | [产品A名称] | [详细描述产品A的关键功能、技术参数和独特卖点] | [说明产品A主要应用于哪些行业领域的哪些具体环节，服务于哪类客户] |
  | [产品B名称] | [详细描述产品B的关键功能、技术参数和独特卖点] | [说明产品B主要应用于哪些行业领域的哪些具体环节，服务于哪类客户] |
  ```

### 3. 核心技术
  请详细阐述公司的核心技术壁垒。具体说明该技术是什么（例如：某种特定算法、新材料配方、专有工艺流程等）。进一步解释该技术是如何实现的，其关键的创新点在哪里？这项技术是否已获得专利保护？相比于竞争对手的技术，它具备哪些显著的优势（如性能、成本、效率等）？

### 4. 核心人员与团队
  **必须使用Markdown表格** 来展示核心团队成员，第一行必须是创始人/CEO，随后是其他核心高管。
  ```
  | 姓名 | 职位 | 学历背景与工作履历 | 相关荣誉与成就 |
  |---|---|---|---|
  | [创始人姓名] | [CEO/董事长] | [毕业院校、专业、过往在知名公司的关键任职经历] | [获得的行业重大奖项、入选的人才计划等] |
  | [核心成员A] | [CTO/技术总监] | [毕业院校、专业、在技术领域的关键项目经验] | [发表的重要论文、拥有的核心专利等] |
  ```
  **表格结束后**，如果文档中有提及团队的整体规模、部门构成（如研发、销售、市场人员比例），请务- 务必在表格结束之后分段用一句话总结。如果未提及，就不用写。

### 5. 市场与竞争
  **市场规模**：(此为独立段落) 如果文档中提供了市场规模的数据（如TAM, SAM, SOM），请用**完整的文字段落**进行描述，说明其引用的数据来源和年份。不要使用列表。
  **竞争格局**：(此为独立段落) 根据文档内容，并根据公司产品或者应用场景分类，清晰列出公司的主要竞争对手。如果可能，将竞争对手分为直接竞争者和间接竞争者。此处仅需列出名字，不要求进行详细的竞品分析。

### 6. 财务与融资
  **财务部分**: **必须使用Markdown表格** 展示近三年及最近一期的核心财务数据。**如果文档中没有提供连续三年的数据或任何财务数据，必须如实说明，严禁杜撰。**
  ```
  | 指标 (单位：万元) | 近期]| 2024年]| 2023年]|
  |---|---|---|---|
  | 营业收入 | | | |
  | 毛利率 (%) | | | |
  | 净利润 | | | |
  | 净利润率 (%) | | | |
  ```
  **融资部分**: **必须使用Markdown表格** 展示公司的历史融资情况。对于投资机构，**仅在文档明确提到投资金额时**，才将金额以括号形式追加在机构名称后。
  ```
  | 轮次 | 本轮融资金额(万元) | 投后估值(万元) | 投资机构 |
  |---|---|---|---|
  | [例如：天使轮] | | | |
  | [例如：A轮] | | | |
  ```"""
    
    # --- 并购模板 (Merger & Acquisition Template) - 指令仍需进一步优化 ---
    merger_summary_template = """### 1. 交易背景与逻辑
  请从宏观行业趋势（例如行业整合、技术变革）、市场竞争格局出发，深入阐述本次交易的根本驱动因素和战略意图。明确说明买方希望通过此次并购实现什么核心目标（例如：获取关键技术、进入新市场、扩大产能、增强供应链控制力、消除竞争、实现客户资源协同等）。如果文档提及，请说明交易的催化事件（例如：标的公司主动寻求出售、股东退出需求等）。

### 2. 买方与标的公司介绍
  请分别对买方和标的公司进行简洁但关键的介绍。对于买方，简要说明其主营业务、行业地位及本次并购前的规模体量。对于标的公司，清晰描述其核心业务、主要产品/服务、技术壁垒、市场地位及其独特的吸引力。此部分的重点在于精准地回答“买方究竟看中了标的公司的什么核心价值？”

### 3. 交易对价情况
  请详细说明本次交易的对价结构。**用加粗标题引导，分段阐述**，例如：**交易总金额**：[段落说明]。**支付方式**：[段落说明]。**估值倍数**：[段落说明]。**溢价情况**：[段落说明]。**特殊条款**：[段落说明]。

### 4. 核心风险
  请从商业、财务、法律和整合等多个维度，全面识别并阐述本次交易面临的核心风险。**用加粗标题引导，分段阐述**，例如：**整合风险**：[段落说明]。**财务风险**：[段落说明]。**商业风险**：[段落说明]。**审批风险**：[段落说明]。

### 5. 整合计划
  请根据文档内容，概述交易完成后的整合计划（Post-Merger Integration, PMI）。**用加粗标题引导，分段阐述**，例如：**协同效应实现路径**：[段落说明]。**治理与人事安排**：[段落说明]。**业务与技术整合**：[段落说明]。如果文档未提及详细计划，必须明确指出。

### 6. 估值与财务
  请总结标的公司的核心财务数据，并说明本次交易的估值方法。
  **财务部分**: **必须使用Markdown表格** 展示标的公司近三年及一期的关键财务指标（营收、毛利、EBITDA、净利润等）。若数据不全，有几年就展示几年。
  ```
  | 指标 (万元) | 近期 | 2024年 | 2023年 |
  |---|---|---|---|
  | 营收 | | | |
  | 毛利率 | | | |
  | 净利润 | | | |
  | 净利润率 | | | |
  ```
  **估值部分**: **表格下方必须用文字段落** 清晰阐述本次交易所采用的主要估值方法（如可比公司分析法(CCA)、可比交易分析法(CTA)、现金流折现法(DCF)等）。如果文档提供了估值过程中的关键假设（如WACC、永续增长率g等），也请一并说明。"""

    # --- 不良资产模板 (Distressed Asset Template) - 指令仍需进一步优化 ---
    distressed_summary_template = """### 1. 资产包概况
  请清晰、具体地描述此不良资产包的构成。它是一家陷入困境的公司实体、一组不良债权、一处或多处抵押资产（如商业地产、厂房设备），还是其他类型的资产组合？说明该资产/公司的当前运营状态（例如：仍在经营但现金流枯竭、已停产、已进入破产程序等）及其所在的核心行业。

### 2. 债务情况与成因
  这是分析的核心。请深入分析导致资产陷入困境的根本原因和债务的详细情况。**用加粗标题引导，分段阐述**，例如：**困境成因**：[段落说明]。**债务结构**：[段落说明，详细列出总债务、主要债权人、清偿顺序等]。

### 3. 资产核心价值
  在剥离债务问题后，客观评估该资产包中真正有价值的核心部分是什么。分析其有形资产（如土地使用权、房产、关键设备）和无形资产（如品牌、核心专利、客户关系、特许经营权或牌照）的质量、状况和潜在变现能力。关键在于回答：“如果剔除债务，这个资产本身还值钱吗？价值在哪里？”

### 4. 处置与重组方案
  请详细阐述文档中提出的解决方案。方案是寻求整体打包出售、分拆剥离出售核心资产、进行债务重组（如展期、降息、债转股）、引入新的战略投资人进行重整，还是直接进行破产清算？如果提及具体的交易结构（例如在破产法框架下的363出售、预重整方案）或关键条款（例如是否已找到“白马骑士”或意向投资人），请务必详细说明。

### 5. 法律与清算风险
  请识别并说明该项目在法律层面和处置过程中可能面临的主要风险。**用加粗标题引导，分段阐述**，例如：**法律风险**：[段落说明]。**清算/处置风险**：[段落说明]。**程序风险**：[段落说明]。

### 6. 财务与回报分析
  请总结与潜在投资回报相关的财务信息。
  **财务与估值**: **如果文档提及，必须使用Markdown表格** 展示与资产相关的历史财务数据（若有）、资产评估价值或清算价值估算。
  ```
  | 项目 | 金额 | 说明 |
  |---|---|---|
  | 账面总资产 (万元) | | |
  | 总负债 (万元) | | |
  | 资产评估价值 (万元) | | |
  | 清算价值预估 (万元) | | |
  ```
  **回报分析**: **表格下方必须用文字段落** 阐述文档中对投资回报的分析。这应包括预估的投资总成本（收购价+后续投入）、潜在的退出渠道（如重整后出售、分拆出售、IPO）和项目预期的回报率指标（如内部收益率IRR、投资回报倍数MOIC）。如果文档未提供明确的回报分析，必须如实说明“文件中未提及明确的投资回报分析”。"""



    prompt = f"""
请从下面的文档中提取信息，并严格按照指定的JSON格式返回。

**提取字段及要求**:
1.  `name` (项目名称): 字符串。
2.  `property` (属性): 必须是 ["股权", "并购", "不良"] 中的一个。默认为 "股权"。
3.  `status` (状态): 新项目默认为 "尚未开始"。
4.  `industry` (行业): 字符串，识别项目最核心的细分行业。
5.  `description` (项目介绍): 字符串，用一到两句话清晰概括核心业务、主要产品及其应用领域。在描述中，请统一使用“该公司”或“公司”来指代项目主体，不要出现具体的项目名称。结尾不要加句号。
6.  `aiSummary` (AI摘要): 这是一个包含项目深度分析的字符串。请先根据文档内容判断项目属于“股权”、“并购”还是“不良”中的哪一类，然后严格按照该类别对应的模板和以下详细要求进行详尽论述。

    --- **AI摘要内容详细要求 (必须严格遵守)** ---
    **重要原则：所有信息都必须严格基于提供的文档原文。如果原文没有提及相关信息，必须明确指出“文件中未提及”，绝对禁止杜撰！**
    **格式要求：**
    - **标题格式**: 摘要中的六个标题必须使用 Markdown H3 格式 (例如 `### 1. 核心痛点`)。
    - **禁止列表**: **绝对禁止使用任何形式的项目符号或编号列表** (例如 o, -, *, 1., a.)。所有列表形式的信息都必须改写为完整的段落，或由加粗标题引导的独立段落。
    - **表格**: 要求使用表格的地方，必须生成标准的Markdown表格。如果表格内涉及到金额，必须在表格结束后用文字段落说明单位。该文字段落需要与表格几乎没有距离。

    --- **摘要模板选择** ---
    - 如果项目是 **“股权”** 类型，请使用以下模板：
    {equity_summary_template}
    - 如果项目是 **“并购”** 类型，请使用以下模板：
    {merger_summary_template}
    - 如果项目是 **“不良”** 类型，请使用以下模板：
    {distressed_summary_template}

**必须严格按照以下JSON格式输出**:
```json
{{
  "name": "...", "property": "...", "status": "未开始", "industry": "...", "description": "...", "aiSummary": "..."
}}

**文档全文如下**:
---
{full_text[:35000]}
---
"""
    
    # --- 关键修改：将 api_key 传入 ---
    result = call_qwen_api(model=ANALYSIS_MODEL, prompt=prompt, api_key=api_key)

    if result["success"]:
        try:
            content_to_load = result["content"]
            start_index = content_to_load.find('{')
            end_index = content_to_load.rfind('}')
            if start_index != -1 and end_index != -1:
                json_str = content_to_load[start_index:end_index+1]
                structured_data = json.loads(json_str)
                return jsonify(structured_data)
            else:
                raise json.JSONDecodeError("无法在响应中找到有效的JSON对象。", content_to_load, 0)
        except json.JSONDecodeError:
            return jsonify({"error": "AI未能返回有效的JSON格式数据。", "raw_response": result["content"]}), 500
    else:
        return jsonify({"error": result["error"]}), 500

# ==============================================================================
# =================== 板块三 / SECTION 3: API LOGIC & SERVER START =============
# ==============================================================================
# ==============================================================================
# =================== 板块三 / SECTION 3: API LOGIC & SERVER START =============
# ==============================================================================
